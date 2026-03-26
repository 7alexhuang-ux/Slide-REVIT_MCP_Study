#!/usr/bin/env python3
"""
Project Showcase — PPTX 生成器
從 Presentation Markdown 格式生成專業演講簡報

用法:
    python generate_pptx.py input.md [--output output.pptx]

依賴:
    pip install python-pptx PyYAML Pillow
"""

import argparse
import io
import os
import re
import sys
import yaml
from pathlib import Path

# Windows 主控台 UTF-8 輸出支援
if sys.platform == 'win32':
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
    sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

try:
    from PIL import Image
    HAS_PIL = True
except ImportError:
    HAS_PIL = False

# ============================================================
# 常數定義
# ============================================================

# 16:9 投影片尺寸
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# 版面配置（基於 8pt 網格系統 + 建築簡報標準）
MARGIN_LEFT = Inches(1.0)
MARGIN_RIGHT = Inches(1.0)
MARGIN_TOP = Inches(0.75)
MARGIN_BOTTOM = Inches(0.75)
CONTENT_WIDTH = SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT  # 11.333"
TITLE_TOP = Inches(0.75)
TITLE_HEIGHT = Inches(1.0)
CONTENT_TOP = Inches(2.0)
CONTENT_HEIGHT = SLIDE_HEIGHT - CONTENT_TOP - MARGIN_BOTTOM  # 4.75"

# 黃金比例分割（圖文並排用）
GOLDEN_LARGE = 0.618  # 62%
GOLDEN_SMALL = 0.382  # 38%
COLUMN_GAP = Inches(0.5)

# ============================================================
# 主題定義
# ============================================================

THEMES = {
    'minimal-white': {
        'name': '極簡白',
        'bg': '#FFFFFF',
        'bg_accent': '#F2F2F2',
        'title': '#1A1A1A',
        'text': '#333333',
        'accent': '#007ACC',
        'subtitle': '#666666',
        'muted': '#999999',
        'bullet': '#007ACC',
        'font_title': '微軟正黑體',
        'font_body': '微軟正黑體',
        'font_code': 'Consolas',
        'title_size': Pt(36),
        'subtitle_size': Pt(20),
        'body_size': Pt(16),
        'caption_size': Pt(11),
        'decoration': 'bottom-line',
    },
    'architect-dark': {
        'name': '建築深色',
        'bg': '#0F0F0F',
        'bg_accent': '#1E1E1E',
        'title': '#FFFFFF',
        'text': '#CCCCCC',
        'accent': '#B8860B',
        'subtitle': '#999999',
        'muted': '#666666',
        'bullet': '#B8860B',
        'font_title': '微軟正黑體',
        'font_body': '微軟正黑體',
        'font_code': 'Consolas',
        'title_size': Pt(36),
        'subtitle_size': Pt(20),
        'body_size': Pt(16),
        'caption_size': Pt(11),
        'decoration': 'corner-lines',
    },
    'blueprint': {
        'name': '藍圖',
        'bg': '#1B2A4A',
        'bg_accent': '#243556',
        'title': '#FFFFFF',
        'text': '#C8D6E5',
        'accent': '#5DADE2',
        'subtitle': '#85C1E9',
        'muted': '#7F8C8D',
        'bullet': '#5DADE2',
        'font_title': '微軟正黑體',
        'font_body': '微軟正黑體',
        'font_code': 'Consolas',
        'title_size': Pt(36),
        'subtitle_size': Pt(20),
        'body_size': Pt(16),
        'caption_size': Pt(11),
        'decoration': 'grid-accent',
    },
    'concrete': {
        'name': '清水模',
        'bg': '#E8E4E1',
        'bg_accent': '#D5CEC9',
        'title': '#2C2C2C',
        'text': '#4A4A4A',
        'accent': '#8B4513',
        'subtitle': '#6B6B6B',
        'muted': '#8B8B8B',
        'bullet': '#8B4513',
        'font_title': '微軟正黑體',
        'font_body': '微軟正黑體',
        'font_code': 'Consolas',
        'title_size': Pt(36),
        'subtitle_size': Pt(20),
        'body_size': Pt(16),
        'caption_size': Pt(11),
        'decoration': 'side-bar',
    },
    'nature': {
        'name': '自然',
        'bg': '#FAFAF5',
        'bg_accent': '#EBE8DF',
        'title': '#2D3436',
        'text': '#3D3D3D',
        'accent': '#27AE60',
        'subtitle': '#636E72',
        'muted': '#95A5A6',
        'bullet': '#27AE60',
        'font_title': '微軟正黑體',
        'font_body': '微軟正黑體',
        'font_code': 'Consolas',
        'title_size': Pt(36),
        'subtitle_size': Pt(20),
        'body_size': Pt(16),
        'caption_size': Pt(11),
        'decoration': 'organic-bar',
    },
}


# ============================================================
# 工具函式
# ============================================================

def hex_to_rgb(hex_str: str) -> RGBColor:
    """將 hex 色碼轉為 RGBColor"""
    h = hex_str.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def set_slide_bg(slide, color_hex: str):
    """設定投影片背景色"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(color_hex)


def add_textbox(slide, left, top, width, height, text, theme,
                font_size=None, color_key='text', bold=False,
                alignment=PP_ALIGN.LEFT, font_key='font_body',
                line_spacing=None):
    """建立文字方塊並回傳 text_frame"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = text
    p.font.name = theme[font_key]
    p.font.size = font_size or theme['body_size']
    p.font.color.rgb = hex_to_rgb(theme[color_key])
    p.font.bold = bold
    p.alignment = alignment
    if line_spacing:
        p.line_spacing = line_spacing

    return tf


def add_bullets(slide, left, top, width, height, bullets, theme,
                font_size=None, color_key='text'):
    """建立條列式文字方塊"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        is_sub = bullet.startswith('  ')
        text = bullet.strip()

        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # 子項目縮排
        if is_sub:
            p.level = 1
            p.font.size = Pt((font_size or theme['body_size']).pt - 2)
        else:
            p.level = 0
            p.font.size = font_size or theme['body_size']

        # 項目符號：主項用短色塊，子項用細線
        bullet_char = '\u2014 ' if not is_sub else '  \u2013 '  # em dash / en dash
        run_bullet = p.add_run()
        run_bullet.text = bullet_char
        run_bullet.font.color.rgb = hex_to_rgb(theme['accent'])
        run_bullet.font.size = p.font.size
        run_bullet.font.name = theme['font_body']

        # 內容文字
        run_text = p.add_run()
        run_text.text = text
        run_text.font.color.rgb = hex_to_rgb(theme[color_key])
        run_text.font.size = p.font.size
        run_text.font.name = theme['font_body']

        # 行距 1.4x + 段落間距
        p.line_spacing = Pt(p.font.size.pt * 1.8)
        p.space_after = Pt(6)

    return tf


def add_placeholder_image(slide, left, top, width, height, description, theme):
    """加入圖片佔位符（灰色方塊 + 描述文字）"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(theme['bg_accent'])
    shape.line.color.rgb = hex_to_rgb(theme['muted'])
    shape.line.width = Pt(1.5)
    # 圓角
    shape.adjustments[0] = 0.02

    tf = shape.text_frame
    tf.word_wrap = True

    # 圖示
    p_icon = tf.paragraphs[0]
    p_icon.alignment = PP_ALIGN.CENTER
    run_icon = p_icon.add_run()
    run_icon.text = '\U0001F4F7'  # 📷
    run_icon.font.size = Pt(40)

    # 空行
    p_spacer = tf.add_paragraph()
    p_spacer.text = ''
    p_spacer.font.size = Pt(8)

    # 說明文字
    p_desc = tf.add_paragraph()
    p_desc.alignment = PP_ALIGN.CENTER
    run_desc = p_desc.add_run()
    run_desc.text = f'請插入：{description}'
    run_desc.font.size = Pt(16)
    run_desc.font.color.rgb = hex_to_rgb(theme['muted'])
    run_desc.font.name = theme['font_body']

    return shape


def add_real_image(slide, left, top, width, height, image_path: Path, theme):
    """嵌入真實圖片，保持比例置中"""
    if not image_path.exists():
        # 找不到圖片，退回佔位符
        return add_placeholder_image(
            slide, left, top, width, height,
            f'找不到: {image_path.name}', theme
        )

    if HAS_PIL:
        with Image.open(image_path) as img:
            img_w, img_h = img.size
        aspect = img_w / img_h
        box_aspect = width / height

        if aspect > box_aspect:
            # 圖片較寬，以寬度為限
            new_w = width
            new_h = int(width / aspect)
        else:
            # 圖片較高，以高度為限
            new_h = height
            new_w = int(height * aspect)

        # 置中偏移
        offset_left = left + (width - new_w) // 2
        offset_top = top + (height - new_h) // 2
        slide.shapes.add_picture(str(image_path), offset_left, offset_top, new_w, new_h)
    else:
        # 沒有 PIL，直接嵌入
        slide.shapes.add_picture(str(image_path), left, top, width, height)


def add_slide_number(slide, num, total, theme):
    """在右下角加入頁碼"""
    txBox = slide.shapes.add_textbox(
        SLIDE_WIDTH - Inches(1.5),
        SLIDE_HEIGHT - Inches(0.5),
        Inches(1.2),
        Inches(0.3),
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f'{num} / {total}'
    p.alignment = PP_ALIGN.RIGHT
    p.font.size = Pt(11)
    p.font.color.rgb = hex_to_rgb(theme['muted'])
    p.font.name = theme['font_body']


# ============================================================
# 裝飾元素
# ============================================================

def add_decoration(slide, theme, slide_type='content'):
    """根據主題加入裝飾元素"""
    deco = theme.get('decoration', 'bottom-line')

    if deco == 'bottom-line':
        # 底部細線
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, SLIDE_HEIGHT - Inches(0.06),
            SLIDE_WIDTH, Inches(0.06),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(theme['accent'])
        shape.line.fill.background()

    elif deco == 'corner-lines':
        # 左上角 L 形裝飾線
        line_len = Inches(1.5)
        line_w = Inches(0.04)
        # 水平線
        h_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.4), Inches(0.4),
            line_len, line_w,
        )
        h_line.fill.solid()
        h_line.fill.fore_color.rgb = hex_to_rgb(theme['accent'])
        h_line.line.fill.background()
        # 垂直線
        v_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.4), Inches(0.4),
            line_w, line_len,
        )
        v_line.fill.solid()
        v_line.fill.fore_color.rgb = hex_to_rgb(theme['accent'])
        v_line.line.fill.background()

        # 右下角 L 形
        h_line2 = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            SLIDE_WIDTH - Inches(0.4) - line_len,
            SLIDE_HEIGHT - Inches(0.4) - line_w,
            line_len, line_w,
        )
        h_line2.fill.solid()
        h_line2.fill.fore_color.rgb = hex_to_rgb(theme['accent'])
        h_line2.line.fill.background()
        v_line2 = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            SLIDE_WIDTH - Inches(0.4) - line_w,
            SLIDE_HEIGHT - Inches(0.4) - line_len,
            line_w, line_len,
        )
        v_line2.fill.solid()
        v_line2.fill.fore_color.rgb = hex_to_rgb(theme['accent'])
        v_line2.line.fill.background()

    elif deco == 'grid-accent':
        # 底部漸層色帶
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, SLIDE_HEIGHT - Inches(0.08),
            SLIDE_WIDTH, Inches(0.08),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = hex_to_rgb(theme['accent'])
        bar.line.fill.background()
        # 右側細線
        side = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            SLIDE_WIDTH - Inches(0.06), 0,
            Inches(0.06), SLIDE_HEIGHT,
        )
        side.fill.solid()
        side.fill.fore_color.rgb = hex_to_rgb(theme['accent'])
        side.line.fill.background()

    elif deco == 'side-bar':
        # 左側粗色帶
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, 0,
            Inches(0.12), SLIDE_HEIGHT,
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = hex_to_rgb(theme['accent'])
        bar.line.fill.background()

    elif deco == 'organic-bar':
        # 底部色帶 + 左側點
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, SLIDE_HEIGHT - Inches(0.07),
            SLIDE_WIDTH, Inches(0.07),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = hex_to_rgb(theme['accent'])
        bar.line.fill.background()

        # 左上角圓點裝飾
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(0.4), Inches(0.4),
            Inches(0.15), Inches(0.15),
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = hex_to_rgb(theme['accent'])
        dot.line.fill.background()


# ============================================================
# Markdown 解析器
# ============================================================

def parse_frontmatter(content: str):
    """解析 YAML frontmatter，回傳 (config_dict, body_str)"""
    if content.startswith('---'):
        parts = content.split('---', 2)
        if len(parts) >= 3:
            try:
                config = yaml.safe_load(parts[1]) or {}
            except yaml.YAMLError:
                config = {}
            return config, parts[2]
    return {}, content


def parse_slides(body: str):
    """將 body 依 --- 分割為投影片列表"""
    # 用獨佔一行的 --- 分割
    raw_slides = re.split(r'\n---\s*\n', body)
    slides = []
    for raw in raw_slides:
        raw = raw.strip()
        if not raw:
            continue
        slide = parse_single_slide(raw)
        if slide:
            slides.append(slide)
    return slides


def parse_single_slide(raw: str) -> dict:
    """解析單一投影片的 Markdown 內容"""
    slide = {
        'type': None,  # 稍後推斷
        'title': '',
        'subtitle': '',
        'bullets': [],
        'images': [],
        'speaker_notes': '',
        'code_blocks': [],
        'raw_text': '',
        'table': None,
    }

    # 提取明確的類型註記
    type_match = re.search(r'<!--\s*(?:type|slide):\s*(\S+)\s*-->', raw)
    if type_match:
        slide['type'] = type_match.group(1)
        raw = re.sub(r'<!--\s*(?:type|slide):\s*\S+\s*-->', '', raw)

    # 提取程式碼區塊（避免干擾其他解析）
    code_blocks = re.findall(r'```(\w*)\n(.*?)```', raw, re.DOTALL)
    slide['code_blocks'] = [{'lang': lang, 'code': code.strip()} for lang, code in code_blocks]
    raw_no_code = re.sub(r'```\w*\n.*?```', '', raw, flags=re.DOTALL)

    # 提取表格
    table_match = re.search(r'(\|.+\|\n\|[-| :]+\|\n(?:\|.+\|\n?)+)', raw_no_code)
    if table_match:
        slide['table'] = table_match.group(1).strip()
        raw_no_code = raw_no_code.replace(table_match.group(0), '')

    lines = raw_no_code.strip().split('\n')
    text_lines = []

    for line in lines:
        stripped = line.strip()

        # H1 → 區段投影片
        if stripped.startswith('# ') and not stripped.startswith('## '):
            slide['title'] = stripped[2:].strip()
            if slide['type'] is None:
                slide['type'] = 'section'

        # H2 → 投影片標題
        elif stripped.startswith('## '):
            slide['title'] = stripped[3:].strip()

        # H3 → 副標題
        elif stripped.startswith('### '):
            slide['subtitle'] = stripped[4:].strip()

        # 條列項目
        elif stripped.startswith('- '):
            slide['bullets'].append(stripped[2:].strip())
        elif stripped.startswith('  - '):
            slide['bullets'].append('  ' + stripped[4:].strip())

        # 圖片
        elif re.match(r'!\[', stripped):
            img_match = re.match(r'!\[(.*?)\]\((.*?)\)', stripped)
            if img_match:
                alt = img_match.group(1)
                src = img_match.group(2)
                if alt.lower().startswith('placeholder:'):
                    slide['images'].append({
                        'type': 'placeholder',
                        'description': alt.split(':', 1)[1].strip(),
                    })
                else:
                    slide['images'].append({
                        'type': 'file',
                        'path': src,
                        'alt': alt,
                    })

        # 講者備註
        elif stripped.startswith('> 講者備註:') or stripped.startswith('> 講者備註：'):
            note_text = re.split(r'[:：]', stripped, maxsplit=1)[1].strip() if re.search(r'[:：]', stripped) else ''
            slide['speaker_notes'] = note_text
        elif stripped.startswith('> Note:'):
            slide['speaker_notes'] = stripped.split(':', 1)[1].strip()

        # 一般文字
        elif stripped and not stripped.startswith('<!--'):
            text_lines.append(stripped)

    slide['raw_text'] = '\n'.join(text_lines)

    # 推斷投影片類型
    if slide['type'] is None:
        if slide['images'] and not slide['bullets']:
            slide['type'] = 'image'
        elif slide['images'] and slide['bullets']:
            slide['type'] = 'split'
        elif slide['table']:
            slide['type'] = 'content'
        elif slide['code_blocks']:
            slide['type'] = 'content'
        else:
            slide['type'] = 'content'

    return slide


# ============================================================
# 投影片生成函式
# ============================================================

def create_cover_slide(prs, config, theme):
    """封面投影片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白版面
    set_slide_bg(slide, theme['bg'])
    add_decoration(slide, theme, 'cover')

    title = config.get('title', '專案展示')
    subtitle = config.get('subtitle', '')
    author = config.get('author', '')
    date = config.get('date', '')

    # 標題 — 置中偏上
    add_textbox(
        slide, MARGIN_LEFT, Inches(2.2),
        CONTENT_WIDTH, Inches(1.2),
        title, theme,
        font_size=Pt(44), color_key='title', bold=True,
        alignment=PP_ALIGN.CENTER, font_key='font_title',
    )

    # 副標題
    if subtitle:
        add_textbox(
            slide, MARGIN_LEFT, Inches(3.5),
            CONTENT_WIDTH, Inches(0.8),
            subtitle, theme,
            font_size=Pt(24), color_key='subtitle',
            alignment=PP_ALIGN.CENTER,
        )

    # 作者 + 日期
    info_text = ' | '.join(filter(None, [author, date]))
    if info_text:
        add_textbox(
            slide, MARGIN_LEFT, Inches(4.8),
            CONTENT_WIDTH, Inches(0.5),
            info_text, theme,
            font_size=Pt(16), color_key='muted',
            alignment=PP_ALIGN.CENTER,
        )

    return slide


def create_section_slide(prs, slide_data, theme):
    """區段標題投影片（分隔不同段落）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, theme['bg'])
    add_decoration(slide, theme, 'section')

    # 大標題置中
    add_textbox(
        slide, MARGIN_LEFT, Inches(2.8),
        CONTENT_WIDTH, Inches(1.5),
        slide_data['title'], theme,
        font_size=Pt(40), color_key='title', bold=True,
        alignment=PP_ALIGN.CENTER, font_key='font_title',
    )

    # 副標題（如果有）
    if slide_data.get('subtitle'):
        add_textbox(
            slide, MARGIN_LEFT, Inches(4.2),
            CONTENT_WIDTH, Inches(0.8),
            slide_data['subtitle'], theme,
            font_size=Pt(22), color_key='subtitle',
            alignment=PP_ALIGN.CENTER,
        )

    # 講者備註
    if slide_data.get('speaker_notes'):
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = slide_data['speaker_notes']

    return slide


def create_content_slide(prs, slide_data, theme, image_mode='placeholder',
                          base_path=None):
    """一般內容投影片（標題 + 條列）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, theme['bg'])
    add_decoration(slide, theme, 'content')

    # 標題
    if slide_data['title']:
        add_textbox(
            slide, MARGIN_LEFT, TITLE_TOP,
            CONTENT_WIDTH, TITLE_HEIGHT,
            slide_data['title'], theme,
            font_size=theme['title_size'], color_key='title', bold=True,
            font_key='font_title',
        )
        # 標題底線
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            MARGIN_LEFT, TITLE_TOP + Inches(1.1),
            Inches(2.5), Inches(0.04),
        )
        line.fill.solid()
        line.fill.fore_color.rgb = hex_to_rgb(theme['accent'])
        line.line.fill.background()

    # 內容區
    content_left = MARGIN_LEFT
    content_width = CONTENT_WIDTH

    # 如果有圖片，縮小文字區域
    if slide_data['images']:
        content_width = CONTENT_WIDTH * 0.55
        img_left = MARGIN_LEFT + content_width + Inches(0.3)
        img_width = CONTENT_WIDTH - content_width - Inches(0.3)
        img_height = CONTENT_HEIGHT - Inches(0.3)

        for img in slide_data['images'][:1]:  # 只取第一張
            if img['type'] == 'placeholder':
                add_placeholder_image(
                    slide, img_left, CONTENT_TOP,
                    img_width, img_height,
                    img['description'], theme,
                )
            elif img['type'] == 'file' and base_path:
                img_path = Path(base_path) / img['path']
                add_real_image(
                    slide, img_left, CONTENT_TOP,
                    img_width, img_height,
                    img_path, theme,
                )

    # 條列式內容
    if slide_data['bullets']:
        add_bullets(
            slide, content_left, CONTENT_TOP,
            content_width, CONTENT_HEIGHT,
            slide_data['bullets'], theme,
        )
    elif slide_data['raw_text']:
        add_textbox(
            slide, content_left, CONTENT_TOP,
            content_width, CONTENT_HEIGHT,
            slide_data['raw_text'], theme,
            color_key='text',
        )

    # 程式碼區塊
    if slide_data['code_blocks']:
        code = slide_data['code_blocks'][0]
        code_top = CONTENT_TOP + Inches(len(slide_data['bullets']) * 0.5 + 0.3)
        # 程式碼背景
        code_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            content_left, code_top,
            content_width, Inches(2.0),
        )
        code_bg.fill.solid()
        code_bg.fill.fore_color.rgb = hex_to_rgb(theme['bg_accent'])
        code_bg.line.fill.background()
        code_bg.adjustments[0] = 0.02
        # 程式碼文字
        add_textbox(
            slide, content_left + Inches(0.3), code_top + Inches(0.2),
            content_width - Inches(0.6), Inches(1.6),
            code['code'], theme,
            font_size=Pt(14), color_key='text',
            font_key='font_code',
        )

    # 表格（簡化版：用文字模擬）
    if slide_data['table']:
        table_top = CONTENT_TOP + Inches(len(slide_data['bullets']) * 0.5 + 0.3)
        _create_table_from_md(slide, slide_data['table'], theme,
                              content_left, table_top, content_width)

    # 講者備註
    if slide_data.get('speaker_notes'):
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = slide_data['speaker_notes']

    return slide


def create_image_slide(prs, slide_data, theme, image_mode='placeholder',
                        base_path=None):
    """大圖投影片（圖片佔滿大部分空間）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, theme['bg'])
    add_decoration(slide, theme, 'image')

    # 標題（較小）
    if slide_data['title']:
        add_textbox(
            slide, MARGIN_LEFT, Inches(0.3),
            CONTENT_WIDTH, Inches(0.8),
            slide_data['title'], theme,
            font_size=Pt(28), color_key='title', bold=True,
            font_key='font_title',
        )

    # 大圖區域
    img_top = Inches(1.3)
    img_height = SLIDE_HEIGHT - img_top - Inches(0.5)
    img_width = CONTENT_WIDTH
    img_left = MARGIN_LEFT

    for img in slide_data['images'][:1]:
        if img['type'] == 'placeholder':
            add_placeholder_image(
                slide, img_left, img_top,
                img_width, img_height,
                img['description'], theme,
            )
        elif img['type'] == 'file' and base_path:
            img_path = Path(base_path) / img['path']
            add_real_image(
                slide, img_left, img_top,
                img_width, img_height,
                img_path, theme,
            )

    # 講者備註
    if slide_data.get('speaker_notes'):
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = slide_data['speaker_notes']

    return slide


def create_split_slide(prs, slide_data, theme, image_mode='placeholder',
                        base_path=None):
    """左文右圖投影片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, theme['bg'])
    add_decoration(slide, theme, 'split')

    # 標題
    if slide_data['title']:
        add_textbox(
            slide, MARGIN_LEFT, TITLE_TOP,
            CONTENT_WIDTH, TITLE_HEIGHT,
            slide_data['title'], theme,
            font_size=theme['title_size'], color_key='title', bold=True,
            font_key='font_title',
        )

    # 左半：文字（黃金比例 38%）
    text_width = int(CONTENT_WIDTH * GOLDEN_SMALL)
    if slide_data['bullets']:
        add_bullets(
            slide, MARGIN_LEFT, CONTENT_TOP,
            text_width, CONTENT_HEIGHT,
            slide_data['bullets'], theme,
        )

    # 右半：圖片（黃金比例 62%）
    img_width = int(CONTENT_WIDTH * GOLDEN_LARGE)
    img_left = MARGIN_LEFT + text_width + COLUMN_GAP
    img_height = CONTENT_HEIGHT - Inches(0.3)

    for img in slide_data['images'][:1]:
        if img['type'] == 'placeholder':
            add_placeholder_image(
                slide, img_left, CONTENT_TOP,
                img_width, img_height,
                img['description'], theme,
            )
        elif img['type'] == 'file' and base_path:
            img_path = Path(base_path) / img['path']
            add_real_image(
                slide, img_left, CONTENT_TOP,
                img_width, img_height,
                img_path, theme,
            )

    # 講者備註
    if slide_data.get('speaker_notes'):
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = slide_data['speaker_notes']

    return slide


def create_comparison_slide(prs, slide_data, theme, base_path=None):
    """對比投影片（兩欄）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, theme['bg'])
    add_decoration(slide, theme, 'comparison')

    # 標題
    if slide_data['title']:
        add_textbox(
            slide, MARGIN_LEFT, TITLE_TOP,
            CONTENT_WIDTH, TITLE_HEIGHT,
            slide_data['title'], theme,
            font_size=theme['title_size'], color_key='title', bold=True,
            font_key='font_title',
        )

    # 如果有表格，解析為兩欄
    if slide_data['table']:
        _create_table_from_md(slide, slide_data['table'], theme,
                              MARGIN_LEFT, CONTENT_TOP, CONTENT_WIDTH)
    elif slide_data['bullets']:
        # 將 bullets 分成兩半
        mid = len(slide_data['bullets']) // 2
        left_bullets = slide_data['bullets'][:mid]
        right_bullets = slide_data['bullets'][mid:]

        half_width = (CONTENT_WIDTH - Inches(0.5)) / 2
        if left_bullets:
            add_bullets(
                slide, MARGIN_LEFT, CONTENT_TOP,
                half_width, CONTENT_HEIGHT,
                left_bullets, theme,
            )
        if right_bullets:
            add_bullets(
                slide, MARGIN_LEFT + half_width + Inches(0.5), CONTENT_TOP,
                half_width, CONTENT_HEIGHT,
                right_bullets, theme,
            )

    # 講者備註
    if slide_data.get('speaker_notes'):
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = slide_data['speaker_notes']

    return slide


def create_quote_slide(prs, slide_data, theme):
    """引言投影片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, theme['bg'])
    add_decoration(slide, theme, 'quote')

    quote_text = slide_data.get('raw_text', slide_data.get('title', ''))

    # 引號裝飾
    add_textbox(
        slide, MARGIN_LEFT, Inches(1.5),
        Inches(1.5), Inches(1.5),
        '\u201C', theme,  # 「
        font_size=Pt(120), color_key='accent',
        alignment=PP_ALIGN.LEFT, font_key='font_title',
    )

    # 引文
    add_textbox(
        slide, MARGIN_LEFT + Inches(1.0), Inches(2.5),
        CONTENT_WIDTH - Inches(2.0), Inches(3.0),
        quote_text, theme,
        font_size=Pt(28), color_key='text',
        alignment=PP_ALIGN.LEFT, font_key='font_body',
        line_spacing=Pt(42),
    )

    # 講者備註
    if slide_data.get('speaker_notes'):
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = slide_data['speaker_notes']

    return slide


def create_demo_slide(prs, slide_data, theme):
    """Demo 投影片（帶「LIVE DEMO」標記）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, theme['bg'])
    add_decoration(slide, theme, 'demo')

    # LIVE DEMO 標記
    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(4.5), Inches(0.4),
        Inches(4.3), Inches(0.8),
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = hex_to_rgb(theme['accent'])
    badge.line.fill.background()
    badge.adjustments[0] = 0.3

    tf_badge = badge.text_frame
    tf_badge.word_wrap = True
    p_badge = tf_badge.paragraphs[0]
    p_badge.alignment = PP_ALIGN.CENTER
    run_badge = p_badge.add_run()
    run_badge.text = '\u25B6  LIVE DEMO'
    run_badge.font.size = Pt(28)
    run_badge.font.bold = True
    run_badge.font.name = theme['font_title']
    # 根據背景深淺決定文字顏色
    bg_hex = theme['accent'].lstrip('#')
    bg_brightness = sum(int(bg_hex[i:i+2], 16) for i in (0, 2, 4)) / 3
    badge_text_color = '#FFFFFF' if bg_brightness < 160 else '#1A1A1A'
    run_badge.font.color.rgb = hex_to_rgb(badge_text_color)

    # 標題
    if slide_data['title']:
        add_textbox(
            slide, MARGIN_LEFT, Inches(1.8),
            CONTENT_WIDTH, Inches(0.8),
            slide_data['title'], theme,
            font_size=Pt(32), color_key='title', bold=True,
            font_key='font_title',
        )

    # Demo 步驟
    if slide_data['bullets']:
        # 用編號而非圓點
        numbered_bullets = []
        for i, b in enumerate(slide_data['bullets'], 1):
            if b.startswith('  '):
                numbered_bullets.append(b)  # 子項目保持原樣
            else:
                numbered_bullets.append(f'Step {i}：{b}')
        add_bullets(
            slide, MARGIN_LEFT, Inches(2.8),
            CONTENT_WIDTH, Inches(4.0),
            numbered_bullets, theme,
        )

    # 講者備註
    if slide_data.get('speaker_notes'):
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = slide_data['speaker_notes']

    return slide


def create_closing_slide(prs, slide_data, theme, config):
    """結語投影片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, theme['bg'])
    add_decoration(slide, theme, 'closing')

    # 核心 Takeaway
    takeaway = slide_data.get('title', '') or slide_data.get('raw_text', '')
    if takeaway:
        add_textbox(
            slide, MARGIN_LEFT, Inches(2.0),
            CONTENT_WIDTH, Inches(1.5),
            takeaway, theme,
            font_size=Pt(36), color_key='title', bold=True,
            alignment=PP_ALIGN.CENTER, font_key='font_title',
        )

    # 條列要點（如果有）
    if slide_data['bullets']:
        add_bullets(
            slide, Inches(3.0), Inches(3.5),
            Inches(7.0), Inches(2.5),
            slide_data['bullets'], theme,
            font_size=Pt(22),
        )

    # 聯絡資訊
    author = config.get('author', '')
    contact = config.get('contact', '')
    info = ' | '.join(filter(None, [author, contact]))
    if info:
        add_textbox(
            slide, MARGIN_LEFT, Inches(5.8),
            CONTENT_WIDTH, Inches(0.5),
            info, theme,
            font_size=Pt(16), color_key='muted',
            alignment=PP_ALIGN.CENTER,
        )

    # 講者備註
    if slide_data.get('speaker_notes'):
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = slide_data['speaker_notes']

    return slide


# ============================================================
# 表格輔助函式
# ============================================================

def _create_table_from_md(slide, table_md, theme, left, top, width):
    """從 Markdown 表格語法建立 PPTX 表格"""
    lines = [l.strip() for l in table_md.strip().split('\n') if l.strip()]
    if len(lines) < 2:
        return

    # 解析表頭
    headers = [c.strip() for c in lines[0].strip('|').split('|')]
    # 跳過分隔線（第二行）
    rows = []
    for line in lines[2:]:
        cells = [c.strip() for c in line.strip('|').split('|')]
        rows.append(cells)

    n_cols = len(headers)
    n_rows = len(rows) + 1  # +1 for header

    # 建立表格
    col_width = int(width / n_cols)
    row_height = Inches(0.55)
    table_shape = slide.shapes.add_table(
        n_rows, n_cols,
        left, top, width, row_height * n_rows,
    )
    table = table_shape.table

    # 設定表頭
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.name = theme['font_body']
            paragraph.font.size = Pt(16)
            paragraph.font.bold = True
            paragraph.font.color.rgb = hex_to_rgb(theme['title'])
            paragraph.alignment = PP_ALIGN.CENTER
        # 表頭背景
        cell.fill.solid()
        cell.fill.fore_color.rgb = hex_to_rgb(theme['bg_accent'])

    # 設定資料列
    for r, row in enumerate(rows):
        for c, value in enumerate(row):
            if c >= n_cols:
                break
            cell = table.cell(r + 1, c)
            cell.text = value
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.name = theme['font_body']
                paragraph.font.size = Pt(15)
                paragraph.font.color.rgb = hex_to_rgb(theme['text'])
                paragraph.alignment = PP_ALIGN.CENTER


# ============================================================
# 主流程
# ============================================================

def generate_pptx(md_path: str, output_path: str = None, theme_override: str = None):
    """主函式：從 Presentation MD 生成 PPTX"""
    md_file = Path(md_path)
    if not md_file.exists():
        print(f'錯誤：找不到檔案 {md_path}', file=sys.stderr)
        sys.exit(1)

    content = md_file.read_text(encoding='utf-8')
    config, body = parse_frontmatter(content)

    # 設定（CLI 覆蓋 frontmatter）
    theme_id = theme_override or config.get('theme', 'minimal-white')
    if theme_id not in THEMES:
        print(f'警告：未知主題 "{theme_id}"，使用 minimal-white', file=sys.stderr)
        theme_id = 'minimal-white'
    theme = THEMES[theme_id]

    image_mode = config.get('image_mode', 'placeholder')
    base_path = config.get('project_path', str(md_file.parent))

    # 解析投影片
    slides_data = parse_slides(body)
    if not slides_data:
        print('錯誤：未找到任何投影片內容', file=sys.stderr)
        sys.exit(1)

    # 建立簡報
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # 封面
    create_cover_slide(prs, config, theme)

    # 逐頁生成
    slide_creators = {
        'section': lambda sd: create_section_slide(prs, sd, theme),
        'content': lambda sd: create_content_slide(prs, sd, theme, image_mode, base_path),
        'image': lambda sd: create_image_slide(prs, sd, theme, image_mode, base_path),
        'split': lambda sd: create_split_slide(prs, sd, theme, image_mode, base_path),
        'comparison': lambda sd: create_comparison_slide(prs, sd, theme, base_path),
        'quote': lambda sd: create_quote_slide(prs, sd, theme),
        'demo': lambda sd: create_demo_slide(prs, sd, theme),
        'closing': lambda sd: create_closing_slide(prs, sd, theme, config),
        'cover': lambda sd: None,  # 封面已在前面建立
    }

    for sd in slides_data:
        creator = slide_creators.get(sd['type'])
        if creator:
            creator(sd)
        else:
            # 預設用 content
            create_content_slide(prs, sd, theme, image_mode, base_path)

    # 加入頁碼（跳過封面）
    total = len(prs.slides)
    for i, slide in enumerate(prs.slides):
        if i == 0:
            continue  # 封面不加頁碼
        add_slide_number(slide, i, total - 1, theme)

    # 輸出
    if output_path is None:
        output_path = str(md_file.with_suffix('.pptx'))
    prs.save(output_path)
    print(f'✅ 已生成 PPTX：{output_path}')
    print(f'   主題：{theme["name"]} ({theme_id})')
    print(f'   投影片數：{len(prs.slides)} 張')
    print(f'   圖片模式：{image_mode}')

    return output_path


# ============================================================
# CLI 入口
# ============================================================

def main():
    parser = argparse.ArgumentParser(
        description='Project Showcase — 從 Presentation MD 生成 PPTX',
    )
    parser.add_argument('input', help='Presentation MD 檔案路徑')
    parser.add_argument('--output', '-o', help='輸出 PPTX 路徑')
    parser.add_argument('--theme', '-t', help=f'主題 ({", ".join(THEMES.keys())})')
    parser.add_argument('--list-themes', action='store_true', help='列出所有可用主題')

    args = parser.parse_args()

    if args.list_themes:
        print('可用主題：')
        for tid, t in THEMES.items():
            print(f'  {tid:20s} {t["name"]}')
        sys.exit(0)

    generate_pptx(args.input, args.output, theme_override=args.theme)


if __name__ == '__main__':
    main()
